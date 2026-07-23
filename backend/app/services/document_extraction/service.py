from __future__ import annotations

import uuid
from pathlib import Path

from .exceptions import ExtractionError
from .schema import (
    BoundingBox,
    ExtractedDocument,
    ExtractedPage,
    ImageObject,
    TableObject,
    TextBlock,
)


def _extract_zip_xml_text(file_path: Path) -> list[str]:
    import zipfile
    import xml.etree.ElementTree as ET
    lines: list[str] = []
    try:
        if not zipfile.is_zipfile(file_path):
            return lines
        with zipfile.ZipFile(file_path, "r") as z:
            xml_files = [f for f in z.namelist() if f.endswith(".xml")]
            xml_files.sort(key=lambda x: (0 if "document.xml" in x else (1 if "slide" in x else 2), x))
            for f_name in xml_files:
                try:
                    xml_content = z.read(f_name)
                    tree = ET.fromstring(xml_content)
                    for node in tree.iter():
                        if node.tag.endswith("t") and node.text:
                            txt = node.text.strip()
                            if txt and len(txt) > 1:
                                lines.append(txt)
                except Exception:
                    continue
    except Exception:
        pass
    return lines


def _extract_printable_strings(file_path: Path) -> list[str]:
    import re
    lines: list[str] = []
    try:
        with open(file_path, "rb") as f:
            raw_bytes = f.read()

        try:
            decoded = raw_bytes.decode("utf-8", errors="ignore")
        except Exception:
            decoded = raw_bytes.decode("latin-1", errors="ignore")

        matches = re.findall(r"[\x20-\x7E\t\r\n]{4,}", decoded)
        for match in matches:
            cleaned = match.strip()
            if len(cleaned) > 5 and not cleaned.startswith("<?xml") and not cleaned.startswith("PK") and not cleaned.startswith("<w:") and not cleaned.startswith("<a:"):
                lines.append(cleaned)
    except Exception:
        pass
    return lines


class DocumentExtractionService:
    """
    Module 1: Document Extraction Engine.

    Converts a PDF, DOCX, PPTX, or TXT document into raw structured information.
    """

    def extract(self, pdf_path: str) -> ExtractedDocument:
        path = Path(pdf_path)
        if not path.exists():
            raise ExtractionError(f"Document file not found: {pdf_path}")

        ext = path.suffix.lower()
        if ext in {".docx", ".doc"}:
            return self._extract_docx(path)
        elif ext in {".pptx", ".ppt"}:
            return self._extract_pptx(path)
        elif ext in {".txt", ".md"}:
            return self._extract_txt(path)
        else:
            return self._extract_pdf(path)

    def _extract_pdf(self, path: Path) -> ExtractedDocument:
        try:
            import fitz  # PyMuPDF
        except ImportError as e:
            raise ExtractionError(
                "PyMuPDF is required for PDF extraction. Install with `pip install pymupdf`."
            ) from e

        try:
            doc = fitz.open(str(path))
        except Exception as e:
            return self._extract_fallback(path, f"PyMuPDF failed: {e}")

        pages: list[ExtractedPage] = []
        try:
            for page_index in range(doc.page_count):
                page = doc.load_page(page_index)
                pages.append(self._extract_page(page, page_index + 1))
        finally:
            metadata = dict(doc.metadata or {})
            doc.close()

        return ExtractedDocument(
            source_path=str(path),
            page_count=len(pages),
            metadata=metadata,
            pages=pages,
        )

    def _extract_docx(self, docx_path: Path) -> ExtractedDocument:
        try:
            import docx
            doc = docx.Document(str(docx_path))

            pages: list[ExtractedPage] = []
            current_blocks: list[TextBlock] = []
            page_number = 1
            word_count = 0

            for para in doc.paragraphs:
                text = para.text.strip()
                if not text:
                    continue

                style_name = para.style.name.lower() if para.style else ""
                bold = any(run.bold for run in para.runs if run.bold is not None)
                italic = any(run.italic for run in para.runs if run.italic is not None)

                font_size = 12.0
                if "heading 1" in style_name or "title" in style_name:
                    font_size = 20.0
                    bold = True
                elif "heading 2" in style_name:
                    font_size = 16.0
                    bold = True
                elif "heading 3" in style_name:
                    font_size = 14.0
                    bold = True

                current_blocks.append(TextBlock(
                    text=text,
                    font=para.style.name if para.style else "Normal",
                    font_size=font_size,
                    bold=bold,
                    italic=italic,
                    bbox=BoundingBox(0, 0, 100, 100),
                ))

                word_count += len(text.split())
                if word_count >= 400:
                    pages.append(ExtractedPage(page_number=page_number, width=612.0, height=792.0, blocks=current_blocks, images=[], tables=[]))
                    page_number += 1
                    current_blocks = []
                    word_count = 0

            for table in doc.tables:
                table_text_lines = []
                for row in table.rows:
                    row_text = [cell.text.strip() for cell in row.cells if cell.text.strip()]
                    if row_text:
                        table_text_lines.append(" | ".join(row_text))
                if table_text_lines:
                    current_blocks.append(TextBlock(
                        text="\n".join(table_text_lines),
                        font="TableText",
                        font_size=11.0,
                        bold=False,
                        italic=False,
                        bbox=BoundingBox(0, 0, 100, 100),
                    ))

            if current_blocks or not pages:
                pages.append(ExtractedPage(page_number=page_number, width=612.0, height=792.0, blocks=current_blocks, images=[], tables=[]))

            return ExtractedDocument(
                source_path=str(docx_path),
                page_count=len(pages),
                metadata={"title": docx_path.stem},
                pages=pages,
            )
        except Exception as e:
            return self._extract_fallback(docx_path, f"python-docx error: {e}")

    def _extract_pptx(self, pptx_path: Path) -> ExtractedDocument:
        try:
            import pptx
            prs = pptx.Presentation(str(pptx_path))
            pages: list[ExtractedPage] = []

            for slide_index, slide in enumerate(prs.slides):
                blocks: list[TextBlock] = []
                for shape in slide.shapes:
                    if not shape.has_text_frame:
                        continue
                    for paragraph in shape.text_frame.paragraphs:
                        text = paragraph.text.strip()
                        if not text:
                            continue

                        bold = any(run.font.bold for run in paragraph.runs if run.font and run.font.bold is not None)
                        italic = any(run.font.italic for run in paragraph.runs if run.font and run.font.italic is not None)
                        font_size = 14.0
                        if paragraph.runs and paragraph.runs[0].font and paragraph.runs[0].font.size:
                            try:
                                font_size = float(paragraph.runs[0].font.size.pt)
                            except Exception:
                                font_size = 14.0

                        if shape == slide.shapes.title:
                            font_size = 22.0
                            bold = True

                        blocks.append(TextBlock(
                            text=text,
                            font="SlideText",
                            font_size=font_size,
                            bold=bold,
                            italic=italic,
                            bbox=BoundingBox(0, 0, 100, 100),
                        ))

                pages.append(ExtractedPage(
                    page_number=slide_index + 1,
                    width=960.0,
                    height=540.0,
                    blocks=blocks,
                    images=[],
                    tables=[],
                ))

            if not pages:
                pages.append(ExtractedPage(page_number=1, width=960.0, height=540.0, blocks=[], images=[], tables=[]))

            return ExtractedDocument(
                source_path=str(pptx_path),
                page_count=len(pages),
                metadata={"title": pptx_path.stem},
                pages=pages,
            )
        except Exception as e:
            return self._extract_fallback(pptx_path, f"python-pptx error: {e}")

    def _extract_fallback(self, file_path: Path, reason: str = "") -> ExtractedDocument:
        lines = _extract_zip_xml_text(file_path)
        if not lines:
            lines = _extract_printable_strings(file_path)

        if not lines:
            raise ExtractionError(f"Could not extract readable text from document ({file_path.name}). {reason}")

        pages: list[ExtractedPage] = []
        current_blocks: list[TextBlock] = []
        page_number = 1

        for line in lines:
            trimmed = line.strip()
            if not trimmed:
                continue

            current_blocks.append(TextBlock(
                text=trimmed,
                font="FallbackText",
                font_size=12.0,
                bold=False,
                italic=False,
                bbox=BoundingBox(0, 0, 100, 100),
            ))

            if len(current_blocks) >= 40:
                pages.append(ExtractedPage(page_number=page_number, width=612.0, height=792.0, blocks=current_blocks, images=[], tables=[]))
                page_number += 1
                current_blocks = []

        if current_blocks or not pages:
            pages.append(ExtractedPage(page_number=page_number, width=612.0, height=792.0, blocks=current_blocks, images=[], tables=[]))

        return ExtractedDocument(
            source_path=str(file_path),
            page_count=len(pages),
            metadata={"title": file_path.stem, "extraction_note": reason},
            pages=pages,
        )

    def _extract_txt(self, txt_path: Path) -> ExtractedDocument:
        try:
            with open(txt_path, "r", encoding="utf-8", errors="ignore") as f:
                content = f.read()
        except Exception as e:
            raise ExtractionError(f"Could not read text file: {e}") from e

        lines = content.splitlines()
        pages: list[ExtractedPage] = []
        current_blocks: list[TextBlock] = []
        page_number = 1

        for line in lines:
            trimmed = line.strip()
            if not trimmed:
                continue

            is_heading = trimmed.startswith("#")
            font_size = 18.0 if is_heading else 12.0
            bold = is_heading

            current_blocks.append(TextBlock(
                text=trimmed,
                font="Markdown" if txt_path.suffix.lower() == ".md" else "PlainText",
                font_size=font_size,
                bold=bold,
                italic=False,
                bbox=BoundingBox(0, 0, 100, 100),
            ))

            if len(current_blocks) >= 40:
                pages.append(ExtractedPage(page_number=page_number, width=612.0, height=792.0, blocks=current_blocks, images=[], tables=[]))
                page_number += 1
                current_blocks = []

        if current_blocks or not pages:
            pages.append(ExtractedPage(page_number=page_number, width=612.0, height=792.0, blocks=current_blocks, images=[], tables=[]))

        return ExtractedDocument(
            source_path=str(txt_path),
            page_count=len(pages),
            metadata={"title": txt_path.stem},
            pages=pages,
        )

    def _extract_page(self, page, page_number: int) -> ExtractedPage:
        import fitz

        rect = page.rect
        text_dict = page.get_text("dict")
        blocks: list[TextBlock] = []

        for block in text_dict.get("blocks", []):
            if block.get("type") != 0:  # 0 = text block, 1 = image block
                continue
            for line in block.get("lines", []):
                for span in line.get("spans", []):
                    text = span.get("text", "")
                    if not text.strip():
                        continue
                    flags = span.get("flags", 0)
                    font_name = span.get("font", "")
                    bold = bool(flags & (1 << 4)) or "bold" in font_name.lower()
                    italic = bool(flags & (1 << 1)) or "italic" in font_name.lower() or "oblique" in font_name.lower()
                    bx = span.get("bbox", [0, 0, 0, 0])
                    blocks.append(TextBlock(
                        text=text,
                        font=font_name,
                        font_size=round(span.get("size", 0), 2),
                        bold=bold,
                        italic=italic,
                        bbox=BoundingBox(*bx),
                    ))

        images: list[ImageObject] = []
        for img in page.get_images(full=True):
            xref = img[0]
            try:
                rects = page.get_image_rects(xref)
                bbox = rects[0] if rects else fitz.Rect(0, 0, 0, 0)
                base = page.parent.extract_image(xref)
                images.append(ImageObject(
                    image_id=f"p{page_number}_img{xref}",
                    bbox=BoundingBox(bbox.x0, bbox.y0, bbox.x1, bbox.y1),
                    width=base.get("width", 0),
                    height=base.get("height", 0),
                    ext=base.get("ext", "png"),
                ))
            except Exception:
                continue

        tables = self._extract_tables(page, page_number)

        return ExtractedPage(
            page_number=page_number,
            width=rect.width,
            height=rect.height,
            blocks=blocks,
            images=images,
            tables=tables,
        )

    def _extract_tables(self, page, page_number: int) -> list[TableObject]:
        """Best-effort table extraction via pdfplumber. Silently skipped if unavailable."""
        tables: list[TableObject] = []
        try:
            import pdfplumber
        except ImportError:
            return tables

        try:
            with pdfplumber.open(page.parent.name) as pdf:
                pl_page = pdf.pages[page_number - 1]
                for i, table in enumerate(pl_page.find_tables()):
                    rows = table.extract() or []
                    x0, top, x1, bottom = table.bbox
                    tables.append(TableObject(
                        table_id=f"p{page_number}_t{i}_{uuid.uuid4().hex[:6]}",
                        bbox=BoundingBox(x0, top, x1, bottom),
                        rows=[[cell or "" for cell in row] for row in rows],
                    ))
        except Exception:
            return tables
        return tables
